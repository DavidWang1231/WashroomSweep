#!/usr/bin/env python3
"""WashroomSweep — 8-slide deck as an editable .pptx (16:9)."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

GROUND = RGBColor(0x0A, 0x10, 0x17)
PANEL  = RGBColor(0x12, 0x1D, 0x29)
RULE   = RGBColor(0x24, 0x34, 0x3F)
INK    = RGBColor(0xD5, 0xE0, 0xEB)
DIM    = RGBColor(0x8A, 0x9F, 0xB4)
FAINT  = RGBColor(0x5D, 0x72, 0x87)
ACCENT = RGBColor(0x35, 0xD0, 0xD8)
ALERT  = RGBColor(0xFF, 0x5D, 0x5D)
GOOD   = RGBColor(0x4A, 0xD9, 0x91)

SANS, MONO = "Helvetica Neue", "Menlo"
W, H = Inches(13.333), Inches(7.5)
M = Inches(0.9)
CW = W - 2 * M

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide(note=""):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = GROUND
    if note:
        s.notes_slide.notes_text_frame.text = note
    return s


def rect(s, x, y, w, h, fill=None, line=None, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if rounded:
        shp.adjustments[0] = 0.05
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for t, sz, col, bold, fnt, after in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.line_spacing = sp; p.space_after = Pt(after)
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.color.rgb = col
        r.font.bold = bold; r.font.name = fnt
    return tb


def eyebrow(s, t, y=Inches(0.8)):
    txt(s, M, y, CW, Inches(0.3), [(t.upper(), 12, ACCENT, True, SANS, 0)])


def title(s, t, y=Inches(1.2), size=38):
    txt(s, M, y, CW, Inches(1.5), [(t, size, INK, True, SANS, 0)], sp=0.95)


def cards(s, items, y, h, cols=None, accents=None):
    cols = cols or len(items)
    gap = Inches(0.28)
    w = int((CW - gap * (cols - 1)) / cols)
    for i, (head, bodytxt) in enumerate(items):
        x = M + i * (w + gap)
        rect(s, x, y, w, h, fill=PANEL, line=RULE)
        hcol = INK
        if accents and accents[i]:
            bar = rect(s, x, y, Pt(4), h, fill=accents[i], rounded=False)
            bar.line.fill.background()
            hcol = accents[i]
        pad = Inches(0.28)
        txt(s, x + pad, y + pad, w - 2 * pad, h - 2 * pad,
            [(head, 16, hcol, True, SANS, 8),
             (bodytxt, 12.5, DIM, False, SANS, 0)], sp=1.2)


def bigs(s, items, y, h=Inches(1.8)):
    gap = Inches(0.28)
    n = len(items)
    w = int((CW - gap * (n - 1)) / n)
    for i, (num, lbl, sub, col) in enumerate(items):
        x = M + i * (w + gap)
        rect(s, x, y, w, h, fill=PANEL, line=RULE)
        bar = rect(s, x, y, w, Pt(3.5), fill=col, rounded=False)
        bar.line.fill.background()
        pad = Inches(0.28)
        runs = [(num, 34, col, True, SANS, 6),
                (lbl.upper(), 11, DIM, True, SANS, 5)]
        if sub:
            runs.append((sub, 10.5, FAINT, False, SANS, 0))
        txt(s, x + pad, y + pad, w - 2 * pad, h - 2 * pad, runs, sp=1.15)


def body(s, t, y, size=16, col=DIM, w=None):
    txt(s, M, y, w or Inches(10.5), Inches(1.3), [(t, size, col, False, SANS, 0)], sp=1.35)


# ─────────────────────────────────── 1
s = slide("开场。一句话说清楚:这是一个进门前的扫查器 —— 在你走进洗手间之前,"
          "告诉你这里是不是有摄像头正在把画面传出去。\n\n"
          "不要一上来讲技术,先让评委理解这是一个每个人都可能遇到的场景。\n\n"
          "队友可以在这里介绍团队。")
eyebrow(s, "Wireless for Humanity · Student Design Competition", Inches(1.9))
title(s, "WashroomSweep", Inches(2.4), size=60)
body(s, "Before you walk into a washroom, know whether a camera in there "
        "is streaming right now.", Inches(3.9), size=22, col=INK)
txt(s, M, Inches(5.1), CW, Inches(0.4),
    [("University of Waterloo  ·  2-person team", 13, FAINT, False, MONO, 0)])

# ─────────────────────────────────── 2
s = slide("背景。洗手间是最典型的场景:你无法预先检查、安装者有充足时间、受害者毫无察觉。\n\n"
          "Airbnb 2024年4月30日起全面禁止室内摄像头 —— 真实政策,可以直接引用。\n\n"
          "市面产品:几十块的 RF 探测器只会说''这里有信号'',但现代建筑里永远有信号,"
          "响个不停等于没响。专业设备有效但要几千块还需要培训。\n\n"
          "要加具体统计数字的话自己查来源,别凭印象写。")
eyebrow(s, "Background")
title(s, "Cheap cameras, and no good way to check")
cards(s, [
    ("The problem is real",
     "A streaming camera now fits inside a USB charger or a clothes hook and costs "
     "less than a meal. Airbnb banned indoor cameras outright in April 2024; South "
     "Korea runs municipal inspection teams for public washrooms."),
    ("Existing tools don't answer the question",
     "RF detectors ($20–60) beep at any radio signal — always true indoors, so "
     "useless. Lens finders need the right angle. Professional gear works but costs "
     "$1k–10k and needs training."),
], Inches(2.65), Inches(2.3), cols=2)
body(s, "They tell you a signal exists. We wanted to know: "
        "is something recording here, right now?", Inches(5.35), size=18, col=INK)

# ─────────────────────────────────── 3
s = slide("我们做了什么。一块几十块的 ESP32 被动监听空气,加一百行 Python 做判断。\n\n"
          "强调''被动'':不加入网络、不解密任何东西,只读 802.11 帧头。"
          "这既是技术特点也是合规特点。\n\n"
          "方向判断用 ToDS/FromDS 位,不需要知道路由器是谁 —— "
          "真实的偷拍摄像头在你没见过的网络上。")
eyebrow(s, "What we built")
title(s, "One ESP32, listening passively")
cards(s, [
    ("Sense",
     "An ESP32 counts bytes per device every 200 ms and splits upload from download. "
     "It joins nothing and decrypts nothing — 802.11 headers only."),
    ("Judge",
     "About 100 lines of Python on a laptop looks for the shape of a camera: heavy "
     "one-way upload that reacts when the room changes."),
    ("Report",
     "A live dashboard flags suspects, and states UNKNOWN whenever it cannot see "
     "properly. It never says “safe.”"),
], Inches(2.6), Inches(2.2))
body(s, "Total radio cost: one dev board — no capture card, no SDR.",
     Inches(5.25), size=18, col=INK)

# ─────────────────────────────────── 4
s = slide("演示 1:被动感知。板子什么网络都没连,但已经看到房间里所有设备,"
          "每台的上行和下行分开统计。\n\n"
          "现场就指着滚动的列表讲。26 秒扫完 13 个信道发现 49 台设备 —— 这是真实数据。\n\n"
          "可以说明:列设备本身不难,市面 APP 都能做。我们的价值在于"
          "从这几十台里挑出哪一台在拍这个房间。")
eyebrow(s, "Demo 1")
title(s, "It sees every device in the room")
bigs(s, [("49", "Devices found", "In a 26-second sweep across all 13 channels", ACCENT),
         ("200 ms", "Update rate", "Exact across 59 windows, zero dropouts", ACCENT),
         ("0", "Packets decrypted", "Never joins a network; header fields only", ACCENT)],
     Inches(2.7), Inches(2.0))
body(s, "Listing devices is the easy part — phone apps do it. The value is picking out "
        "which one is watching this room.", Inches(5.2), size=18, col=INK)

# ─────────────────────────────────── 5
s = slide("演示 2:摄像头检测 —— 重头戏。摄像头只上传不下载。"
          "实测 478 KB/s 单向上行,比第二活跃的设备高 170 倍。\n\n"
          "第二个信号讲原理即可:摄像头用可变码率编码,画面变了数据量就变。"
          "关灯开灯做出一个已知节奏,看谁的流量跟着走 —— 跟得上的就是能看到这个房间的。\n\n"
          "必须主动说:为了绕开 ESP32 看不到 WiFi 6 的硬件限制,我们让板子自己建 802.11n 的网。"
          "真实场景摄像头在房间自己的网络上,这是真实约束。不主动说会被评委问倒。")
eyebrow(s, "Demo 2")
title(s, "The camera gives itself away")
cards(s, [
    ("Signal 1 — it only uploads",
     "Phones and laptops mostly pull data down. A streaming camera does the opposite, "
     "and it never stops.\n\n478 KB/s up  ·  0 down\n170× the next busiest device in "
     "the room."),
    ("Signal 2 — it reacts to the room",
     "Cameras use variable-bitrate encoding: change what the camera sees and you "
     "change how much it sends.\n\nFlip the room's light in a known rhythm and see "
     "whose traffic follows. That separates a camera watching this room from someone "
     "uploading a video next door."),
], Inches(2.6), Inches(2.6), cols=2, accents=[ALERT, None])
body(s, "For this demo the board hosts its own 802.11n network, because its radio "
        "cannot demodulate WiFi 6. A real camera sits on the room's own network — "
        "that limit is on the summary slide.", Inches(5.5), size=13, col=FAINT)

# ─────────────────────────────────── 6
s = slide("BLE —— 队友讲。我们也做了蓝牙扫描,但它不参与摄像头判定。\n\n"
          "原因:蓝牙带宽比视频低两三个数量级,摄像头的视频数据在蓝牙上根本不存在。\n\n"
          "队友解决的真实问题:手机蓝牙地址每 15 分钟轮换,表里会堆满不存在的''幽灵设备''。\n\n"
          "这个''评估过并排除了''的回答是加分的 —— 显示做过取舍,不是漏了没做。")
eyebrow(s, "Demo 3 · Bluetooth")
title(s, "We scan BLE too — and keep it out of the verdict", size=34)
cards(s, [
    ("What it adds",
     "A second board lists nearby Bluetooth devices — signal, maker, name, and whether "
     "the address is randomized. Useful context about what else is in the room."),
    ("Why it can't find cameras",
     "Bluetooth carries orders of magnitude less data than live video, so a streaming "
     "camera's payload never appears there. Seeing BLE only tells you a device has a "
     "Bluetooth radio."),
], Inches(2.7), Inches(2.2), cols=2)
body(s, "Hardest part: phones rotate their Bluetooth address every ~15 minutes, so a "
        "naive list fills up with devices that no longer exist. Entries expire fast so "
        "the list reflects what is here now.", Inches(5.3), size=16, col=INK)

# ─────────────────────────────────── 7
s = slide("其他用途。同样的''谁在单向大量上传''这个信号,换个场景就有别的用途。\n\n"
          "讲得轻松一点,展示想象力。但说清楚这些是设想,不是已实现的功能。")
eyebrow(s, "Where else it works")
title(s, "The same signal, other rooms")
cards(s, [
    ("Confidential meetings",
     "Sweep before a sensitive discussion for anything streaming out of the room."),
    ("Exams and classrooms",
     "Sustained one-way upload from a seat looks different from ordinary browsing — "
     "without inspecting any content."),
    ("Rentals and changing rooms",
     "A routine check between guests, instead of waiting for someone to complain."),
], Inches(2.7), Inches(2.0))
body(s, "All of it reuses one primitive: who is sending far more than they receive, "
        "and does it react to this room?", Inches(5.2), size=18, col=INK)
txt(s, M, Inches(6.1), CW, Inches(0.35),
    [("Directions, not shipped features.", 12, FAINT, False, MONO, 0)])

# ─────────────────────────────────── 8
s = slide("总结。我们不声称机制新颖 —— 先行工作已经存在。贡献是低成本实现、"
          "进门扫查场景、用房间自己的灯做刺激,以及量化了方法什么时候失效。\n\n"
          "盲区一定要主动讲:本地存储摄像头和有线摄像头看不见;"
          "5GHz 和 WiFi 6 超出这个射频范围(实测只能恢复 0.03%);RSSI 不等于距离。\n\n"
          "结尾落在诚实上:能说清楚边界在哪,比声称''全做完了''更可信。\n\n"
          "问答准备在 PRESENTATION_NOTES.md 里。")
eyebrow(s, "Summary")
title(s, "What works, and what doesn't")
cards(s, [
    ("Works",
     "Live over real RF, a streaming camera is separated from every other device in "
     "the room by three orders of magnitude — on upload asymmetry alone."),
    ("Blind spots we measured",
     "Cameras recording to an SD card or over cable emit nothing. 5 GHz and WiFi 6 are "
     "outside this radio — we recovered 0.03% of an 802.11ax link. RSSI is not "
     "distance; we don't localize."),
], Inches(2.6), Inches(2.3), cols=2, accents=[GOOD, ALERT])
body(s, "A sweeper that says UNKNOWN when it cannot see is more useful than one that "
        "says SAFE when it guessed.", Inches(5.3), size=18, col=INK)
txt(s, M, Inches(6.2), CW, Inches(0.7),
    [("Prior art exists (DeWiCam, SnoopDog, Lumos). Ours is the low-cost build, the "
      "entry-sweep scenario, and honest numbers for where it fails.", 12, FAINT, False,
      SANS, 4),
     ("github.com/DavidWang1231/WashroomSweep", 12, FAINT, False, MONO, 0)], sp=1.3)

out = sys.argv[1] if len(sys.argv) > 1 else "WashroomSweep.pptx"
prs.save(out)
print(f"wrote {out}")
